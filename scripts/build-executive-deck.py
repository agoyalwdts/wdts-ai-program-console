#!/usr/bin/env python3
"""Generate docs/executive/WDTS-AI-Program-Console-Executive-Brief.pptx (requires python-pptx)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = REPO_ROOT / "docs" / "executive" / "WDTS-AI-Program-Console-Executive-Brief.pptx"


def _bullets(slide, title: str, lines: list[str]) -> None:
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def main() -> None:
    prs = Presentation()

    # Slide 1 — Title
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "WDTS AI Program Console"
    s0.placeholders[1].text = (
        "One place to see cost, seats, and guardrail-program health\n\n"
        "Briefing for Steering · Security · FinOps"
    )

    layout = prs.slide_layouts[1]

    _bullets(
        prs.slides.add_slide(layout),
        "Why this matters",
        [
            "Generative AI is a line item and an operational risk.",
            "Multiple approved vendors — each with its own admin UI and billing shape.",
            "Money, seats, and policy drift without a shared operating picture.",
            "Leadership needs one trusted view aligned to WDTS policy.",
            "Proposal: a governed internal console — the operational cockpit, not a new policy silo.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "The problem today",
        [
            "Fragmentation: hard to answer “Are we on budget?” across Cursor, OpenAI, Copilot, Claude.",
            "Vendor dashboards don’t encode WDTS planning envelopes (pooled credits, prepaid commit, caps).",
            "Access to tools is not the same as visibility for FinOps and Steering.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "What we are building",
        [
            "A secure, invite-only dashboard for program-level health.",
            "Five approved products: Cursor · ChatGPT · Codex · Claude.ai · M365 Copilot.",
            "Numbers tied to agreed budgets and contracts.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        'What “good” looks like — Program Health',
        [
            "Planning envelopes: Cursor credit band; OpenAI pool + license + overage planning; Copilot prepaid commit.",
            "Month-to-date reconciled to vendor APIs where live (e.g. Cursor Team Admin, OpenAI org costs, Copilot Graph).",
            "Same definitions Finance and Engineering use in planning.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "Governance model",
        [
            "Policy repo (Git PRs) — rulebook and license inventory.",
            "Microsoft Entra ID — identity and sign-in only.",
            "AI Program Console — who may use this app (invites, roles, permissions).",
            "Closed by default; no prompt bodies stored — metadata and aggregates only.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "Financial posture (snapshot)",
        [
            "Cursor — annual credit envelope; spend vs vendor data when synced.",
            "ChatGPT + Codex — pooled credits, per-seat license line, overage economics.",
            "M365 Copilot — prepaid enterprise commitment; level monthly planning line.",
            "Claude.ai — placeholder until contract final.",
            "One planning total for leadership; FinOps keeps policy repo and code constants aligned.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "Current state",
        [
            "Live: Azure-hosted preview, Entra sign-in, app-level RBAC.",
            "Real integrations where secrets and consent exist; synthetic mode for demos.",
            "Identity reconcile and vendor spend sync (cron + manual).",
            "Maturing: custom domain, full production hardening, write workflows (tiers, reclamation, exceptions) next.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "What we are asking to change",
        [
            "Operating norm: console is the default answer for “where is the AI program?”",
            "FinOps owns program number updates (policy repo + dashboard constants).",
            "Expand invited operators — not “everyone with an Entra account.”",
            "Material policy changes flow through the policy repository.",
            "Prioritize hardening vs workflow features with Security and Steering.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "Risks and mitigations",
        [
            "Console vs policy clarity — charter: policy repo + Entra stay authoritative.",
            "Credentials — Key Vault, least privilege, per-env integration toggles.",
            "Vendor mismatch — API sync and documented reconciliation.",
            "Access creep — invite-only model and auditable admin actions.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "Decisions needed",
        [
            "Confirm governance split (policy repo vs console vs Entra).",
            "Assign FinOps owner for program numbers and review cadence.",
            "Prioritize production hardening vs tier/exception workflows.",
            "Approve production deploy path and named promotion reviewers.",
        ],
    )

    _bullets(
        prs.slides.add_slide(layout),
        "Next steps",
        [
            "Steering acknowledgment of charter (deck + policy brief).",
            "FinOps: lock quarterly planning numbers in policy repo and console constants.",
            "Security: confirm integration list and data-handling posture.",
            "Engineering: execute hardening and workflow milestones per prioritization.",
        ],
    )

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
